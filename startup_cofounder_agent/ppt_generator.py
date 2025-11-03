"""Utilities for constructing a pitch deck using python‑pptx.

This module provides a single function, ``create_pitch_deck``, which
takes a list of slide dictionaries (as produced by
``generator.generate_pitch_deck_outline``) and writes a PowerPoint file
to disk. Each slide contains a title and a body consisting of bullet
points separated by newlines.

The function uses the default slide layouts of python‑pptx: the
first slide (layout 0) is a title slide, while subsequent slides use
layout 1 (title and content). Bullet points are added to the content
placeholder.
"""

from __future__ import annotations

from typing import List, Dict
from pptx import Presentation
from pptx.util import Inches, Pt  # type: ignore


def create_pitch_deck(slides: List[Dict[str, str]], output_path: str) -> None:
    """Create a PowerPoint deck from an outline and save it to a file.

    Args:
        slides: A list of dictionaries with keys ``title`` and ``content``.
            The ``content`` should be a string of bullet points separated by
            newlines or semicolons.
        output_path: The filesystem path where the PPTX file should be
            written. The parent directories must already exist.

    Raises:
        ValueError: If no slides are provided.
    """
    if not slides:
        raise ValueError("No slides provided to create_pitch_deck().")

    prs = Presentation()
    for index, slide_data in enumerate(slides):
        title = slide_data.get("title", "")
        content = slide_data.get("content", "")
        # Decide layout: first slide uses title slide layout
        if index == 0:
            slide_layout = prs.slide_layouts[0]
        else:
            slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        # Set the title
        if slide.shapes.title:
            slide.shapes.title.text = title
        # Set the body (content placeholder) if available
        # The second placeholder is typically the body for layout 1
        # For the title slide layout 0 there is a subtitle placeholder
        for shape in slide.placeholders:
            if shape.is_placeholder and shape.placeholder_format.type == 1:  # Body placeholder
                text_frame = shape.text_frame
                # Clear any existing text
                text_frame.clear()
                # Split content into lines; separate on semicolons or newlines
                # Semicolons are common in language model responses
                lines = [line.strip() for line in content.replace(";", "\n").split("\n") if line.strip()]
                if not lines:
                    break
                # Add the first line as the first bullet point
                p = text_frame.paragraphs[0]
                p.text = lines[0]
                # Add subsequent lines as additional bullet points
                for line in lines[1:]:
                    p = text_frame.add_paragraph()
                    p.text = line
                # Apply bullet style
                for paragraph in text_frame.paragraphs:
                    paragraph.level = 0
                break
    # Save the file
    prs.save(output_path)